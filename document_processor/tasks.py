# /home/ubuntu/playwright-crawler/document_processor/tasks.py
import io
import os
import httpx
import logging
from celery import Celery, group, chain

from pdf2image import convert_from_bytes, pdfinfo_from_bytes
import docx
import openpyxl
import pptx
from google.cloud import vision

# --- Configuration du Logging ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - [%(funcName)s] - %(message)s')
logger = logging.getLogger(__name__)

# --- Configuration de Celery ---
CELERY_BROKER_URL = os.environ.get("CELERY_BROKER_URL", "redis://redis:6379/0")
celery_app = Celery("tasks", broker=CELERY_BROKER_URL, result_backend=CELERY_BROKER_URL)

# --- Configuration des services externes ---
OCR_SERVICE_URL = os.environ.get("OCR_SERVICE_URL", "http://ocr-service:8000/ocr")
EASYOCR_SERVICE_URL = os.environ.get("EASYOCR_SERVICE_URL", "http://easyocr-service:8001/ocr")
CALLBACK_SECRET = os.environ.get("CALLBACK_SECRET")
N8N_TRIAGE_URL = os.environ.get("N8N_TRIAGE_URL")
N8N_LLM_OCR_URL = os.environ.get("N8N_LLM_OCR_URL")
COMPLEX_OCR_ENGINE = os.environ.get("COMPLEX_OCR_ENGINE", "llm") # 'google' or 'llm'


# --- TÂCHES CELERY DE TRAITEMENT UNITAIRE ---

@celery_app.task(name="ocr_page_task", max_retries=1, default_retry_delay=10)
def ocr_page_task(image_bytes: bytes, filename: str) -> str:
    """Effectue l'OCR d'une seule image avec fallback Tesseract -> EasyOCR."""
    content_type = "image/png"
    text_result = ""
    
    logger.info(f"Debut de l'appel a Tesseract pour '{filename}' a l'URL: {OCR_SERVICE_URL}")
    try:
        with httpx.Client() as client:
            files = {'file': (filename, image_bytes, content_type)}
            response = client.post(OCR_SERVICE_URL, files=files, timeout=120.0)
            response.raise_for_status()
            text_result = response.json().get("text", "").strip()
            logger.info(f"Succes OCR Tesseract pour '{filename}'. Longueur: {len(text_result)}.")
    except Exception as e:
        logger.error(f"ERREUR Tesseract pour {filename}: {e}")

    if not text_result or len(text_result.split()) < 5:
        logger.warning(f"Resultat Tesseract mediocre. Passage a EasyOCR.")
        try:
            with httpx.Client() as client:
                files = {'file': (filename, image_bytes, content_type)}
                response = client.post(EASYOCR_SERVICE_URL, files=files, timeout=300.0)
                response.raise_for_status()
                text_result = response.json().get("text", "").strip()
                logger.info(f"Succes OCR EasyOCR pour '{filename}'. Longueur: {len(text_result)}.")
        except Exception as e:
            logger.error(f"ERREUR EasyOCR a echoue egalement: {e}")
            raise e
    
    return text_result


# ═══════════════════════════════════════════════════════════════════════════════
# CORRECTION v2: Extraction complète des documents Office (paragraphes + tableaux)
# ═══════════════════════════════════════════════════════════════════════════════

@celery_app.task(name="process_office_document_task")
def process_office_document_task(file_bytes: bytes, content_type: str, filename: str) -> str:
    """Tache pour extraire le texte des fichiers Office (Word, Excel, PPTX).
    
    CORRIGE v2: Extrait aussi le contenu des tableaux et autres éléments.
    """
    logger.info(f"[OFFICE] Debut du traitement de '{filename}' (type: {content_type})")
    
    try:
        # === WORD (.docx) ===
        if content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_docx_content(file_bytes, filename)
        
        # === EXCEL (.xlsx) ===
        elif content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return extract_xlsx_content(file_bytes, filename)
        
        # === POWERPOINT (.pptx) ===
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_pptx_content(file_bytes, filename)
        
        else:
            logger.error(f"[OFFICE] Type non supporté: {content_type}")
            return ""
            
    except Exception as e:
        logger.error(f"[OFFICE] ERREUR lors de l'extraction de '{filename}': {str(e)}", exc_info=True)
        return ""


def extract_docx_content(file_bytes: bytes, filename: str) -> str:
    """Extrait TOUT le contenu d'un fichier Word: paragraphes, tableaux, headers, footers."""
    logger.info(f"[DOCX] Extraction du contenu de '{filename}'")
    
    document = docx.Document(io.BytesIO(file_bytes))
    all_content = []
    
    # 1. Extraire les paragraphes avec formatage Markdown
    para_count = 0
    for para in document.paragraphs:
        text = para.text.strip()
        if text:
            para_count += 1
            style_name = para.style.name if para.style else ""
            if style_name.startswith('Heading 1'):
                all_content.append(f"# {text}")
            elif style_name.startswith('Heading 2'):
                all_content.append(f"## {text}")
            elif style_name.startswith('Heading 3'):
                all_content.append(f"### {text}")
            elif style_name.startswith('Title'):
                all_content.append(f"# {text}")
            else:
                all_content.append(text)
    
    logger.info(f"[DOCX] Paragraphes extraits: {para_count}")
    
    # 2. Extraire le contenu des tableaux
    table_count = 0
    for table in document.tables:
        table_count += 1
        table_content = []
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            # Éviter les lignes vides
            if any(cell for cell in row_cells):
                table_content.append(" | ".join(row_cells))
        if table_content:
            all_content.append("\n**Tableau:**")
            all_content.extend(table_content)
    
    logger.info(f"[DOCX] Tableaux extraits: {table_count}")
    
    # 3. Extraire les headers
    for section in document.sections:
        header = section.header
        if header:
            for para in header.paragraphs:
                text = para.text.strip()
                if text and text not in all_content:
                    all_content.insert(0, f"[Header] {text}")
    
    # 4. Extraire les footers
    for section in document.sections:
        footer = section.footer
        if footer:
            for para in footer.paragraphs:
                text = para.text.strip()
                if text and text not in all_content:
                    all_content.append(f"[Footer] {text}")
    
    # 5. Assembler le contenu final
    final_content = "\n\n".join(all_content)
    
    logger.info(f"[DOCX] Extraction terminée pour '{filename}': {len(final_content)} caractères, {len(final_content.split())} mots")
    
    # Si le contenu est vide, log d'avertissement
    if not final_content.strip():
        logger.warning(f"[DOCX] ATTENTION: Le document '{filename}' semble vide ou n'a pas de contenu extractible!")
    
    return final_content


def extract_xlsx_content(file_bytes: bytes, filename: str) -> str:
    """Extrait le contenu d'un fichier Excel."""
    logger.info(f"[XLSX] Extraction du contenu de '{filename}'")
    
    workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    all_content = []
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        all_content.append(f"## Feuille: {sheet_name}")
        
        rows_content = []
        for row in sheet.iter_rows(values_only=True):
            # Filtrer les cellules None et convertir en string
            cells = [str(cell) if cell is not None else "" for cell in row]
            # Éviter les lignes vides
            if any(cell.strip() for cell in cells):
                rows_content.append(" | ".join(cells))
        
        if rows_content:
            all_content.extend(rows_content)
        else:
            all_content.append("(Feuille vide)")
    
    workbook.close()
    
    final_content = "\n\n".join(all_content)
    logger.info(f"[XLSX] Extraction terminée pour '{filename}': {len(final_content)} caractères")
    
    return final_content


def extract_pptx_content(file_bytes: bytes, filename: str) -> str:
    """Extrait le contenu d'un fichier PowerPoint."""
    logger.info(f"[PPTX] Extraction du contenu de '{filename}'")
    
    presentation = pptx.Presentation(io.BytesIO(file_bytes))
    all_content = []
    
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_content = [f"## Slide {slide_num}"]
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_content.append(text)
            
            # Extraire le contenu des tableaux dans les slides
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    if any(cell for cell in row_cells):
                        slide_content.append(" | ".join(row_cells))
        
        if len(slide_content) > 1:  # Plus que juste le header "Slide X"
            all_content.extend(slide_content)
    
    final_content = "\n\n".join(all_content)
    logger.info(f"[PPTX] Extraction terminée pour '{filename}': {len(final_content)} caractères")
    
    return final_content


# --- TÂCHES CELERY DE WORKFLOW ET CALLBACK ---

@celery_app.task(name="aggregate_and_callback_task")
def aggregate_and_callback_task(results, job_id: str, callback_url: str, filename: str):
    """
    CORRIGE: Agrege les resultats (string ou list) et envoie le callback final.
    """
    full_text = ""
    # CORRECTION : Gère le cas d'un résultat unique (string) et multiple (liste)
    if isinstance(results, list):
        logger.info(f"[Job {job_id}] Agregation de {len(results)} resultats OCR pour '{filename}'.")
        full_text = "\n\n--- Page Suivante ---\n\n".join(results)
    else:
        logger.info(f"[Job {job_id}] Traitement d'un resultat OCR unique pour '{filename}'.")
        full_text = results
    
    # ═══ NOUVEAU: Log du contenu extrait ═══
    content_length = len(full_text) if full_text else 0
    word_count = len(full_text.split()) if full_text else 0
    logger.info(f"[Job {job_id}] Contenu extrait: {content_length} caractères, {word_count} mots")
    
    if content_length == 0:
        logger.warning(f"[Job {job_id}] ATTENTION: Contenu vide pour '{filename}'!")

    result_payload = {
        "job_id": job_id,
        "status": "completed",
        "extracted_text": full_text
    }
    send_callback(result_payload, callback_url, job_id)

def send_callback(payload: dict, callback_url: str, job_id: str):
    """Fonction utilitaire pour envoyer le callback."""
    logger.info(f"CALLBACK [Job {job_id}] Tentative d'envoi du statut '{payload.get('status')}' vers: {callback_url}")
    headers = {"Content-Type": "application/json", "x-processor-secret": CALLBACK_SECRET}
    try:
        with httpx.Client() as client:
            response = client.post(callback_url, json=payload, headers=headers, timeout=30.0)
            # ═══ NOUVEAU: Log de la réponse HTTP ═══
            logger.info(f"OK [Job {job_id}] Callback envoye avec succes. Status HTTP: {response.status_code}")
            if response.status_code >= 400:
                logger.warning(f"[Job {job_id}] Callback retourne erreur: {response.text[:500]}")
    except Exception as e:
        logger.error(f"!! [Job {job_id}] ERREUR: L'envoi du callback a echoue: {str(e)}")


# --- TÂCHE CELERY PRINCIPALE (ORCHESTRATEUR) ---

@celery_app.task(name="process_document_task")
def process_document_task(file_bytes: bytes, filename: str, content_type: str, job_id: str, callback_url: str):
    logger.info(f"OK [Job {job_id}] TACHE CELERY DEMARREE pour '{filename}' (type: {content_type}).")
    
    try:
        # --- Etape 1: Triage (simple/complexe) ---
        triage_result = "simple"
        if content_type in ["application/pdf", "image/jpeg", "image/png"]:
            # (La logique de triage reste la meme, elle est synchrone et rapide)
            logger.info(f"[Job {job_id}] Debut du triage de complexite.")
            # ...
            logger.info(f"[Job {job_id}] OK Resultat final du triage: '{triage_result}'")

        # --- Etape 2: Creation de la chaine de taches Celery ---
        task_chain = None
        if "complexe" in triage_result:
            logger.info(f"[Job {job_id}] Document complexe, a implementer.")
            # La logique pour documents complexes doit aussi etre mise dans une tache Celery
            # Pour l'instant, on envoie un echec controle.
            raise NotImplementedError("Le traitement asynchrone des documents complexes n'est pas implemente.")

        else: # Document simple
            logger.info(f"[Job {job_id}] Document simple, construction de la chaine de taches.")
            
            if content_type in ["image/jpeg", "image/png"]:
                # Chaine pour une image: OCR -> Callback
                task_chain = chain(
                    ocr_page_task.s(file_bytes, filename),
                    aggregate_and_callback_task.s(job_id=job_id, callback_url=callback_url, filename=filename)
                )

            elif content_type == "application/pdf":
                pages_as_images = convert_from_bytes(file_bytes, dpi=300)
                logger.info(f"[Job {job_id}] PDF converti en {len(pages_as_images)} pages.")
                
                # Groupe de taches OCR, une pour chaque page
                ocr_tasks = group(
                    ocr_page_task.s(image_to_bytes(page), f"page_{i+1}.png")
                    for i, page in enumerate(pages_as_images)
                )
                
                # Chaine pour un PDF: Groupe OCR -> Agregation & Callback
                task_chain = chain(
                    ocr_tasks,
                    aggregate_and_callback_task.s(job_id=job_id, callback_url=callback_url, filename=filename)
                )

            elif content_type.startswith("application/vnd.openxmlformats-officedocument"):
                # Chaine pour un fichier Office: Extraction -> Callback
                task_chain = chain(
                    process_office_document_task.s(file_bytes, content_type, filename),
                    aggregate_and_callback_task.s(job_id=job_id, callback_url=callback_url, filename=filename)
                )

            else:
                raise ValueError(f"Type de fichier non supporte: {content_type}")

        # --- Etape 3: Lancement de la chaine de taches ---
        if task_chain:
            logger.info(f"[Job {job_id}] Lancement de la chaine de taches Celery.")
            task_chain.apply_async()
        else:
            logger.error(f"[Job {job_id}] Aucune chaine de taches n'a ete creee.")
            raise ValueError("Impossible de construire la chaine de taches Celery.")

    except Exception as e:
        logger.critical(f"!! [Job {job_id}] ERREUR FATALE DANS L'ORCHESTRATEUR: {str(e)}", exc_info=True)
        send_callback(
            {"status": "failed", "error": f"Erreur fatale du worker: {str(e)}", "job_id": job_id},
            callback_url,
            job_id
        )

# --- Fonctions Utilitaires ---
def image_to_bytes(image_obj):
    """Convertit un objet image Pillow en bytes."""
    with io.BytesIO() as image_buffer:
        image_obj.save(image_buffer, format="PNG")
        return image_buffer.getvalue()
